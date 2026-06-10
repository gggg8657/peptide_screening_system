"""프로그래밍 QA — 텍스트 추출 + 박스 경계/오프슬라이드/겹침 검사 (soffice 없는 환경)."""
import sys
from pptx import Presentation
from pptx.util import Emu

PATH = "/home/dongjukim/Documents/workspace/tmp/SST14-M_scr/_workspace/SSTR2_Architecture.pptx"
SW, SH = 13.333, 7.5  # inches

def emu_in(v): return v / 914400.0

prs = Presentation(PATH)
print(f"slides: {len(prs.slides)}\n")

def boxes(slide):
    out = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        x, y = emu_in(sh.left), emu_in(sh.top)
        w, h = emu_in(sh.width or 0), emu_in(sh.height or 0)
        txt = ""
        if sh.has_text_frame:
            txt = " ".join(p.text for p in sh.text_frame.paragraphs).strip()
        out.append((x, y, w, h, txt, sh.shape_type))
    return out

def overlap(a, b):
    ax, ay, aw, ah = a[:4]; bx, by, bw, bh = b[:4]
    ix = max(0, min(ax+aw, bx+bw) - max(ax, bx))
    iy = max(0, min(ay+ah, by+bh) - max(ay, by))
    return ix * iy

issues = 0
for si, slide in enumerate(prs.slides, 1):
    bs = boxes(slide)
    # off-slide / margin checks
    for (x, y, w, h, txt, st) in bs:
        if x < -0.02 or y < -0.02 or x + w > SW + 0.05 or y + h > SH + 0.05:
            print(f"[S{si}] OFF-SLIDE: ({x:.2f},{y:.2f},{w:.2f},{h:.2f}) '{txt[:30]}'"); issues += 1
    # text-box vs text-box overlap (ignore tiny + decorative)
    txt_boxes = [b for b in bs if b[4]]
    for i in range(len(txt_boxes)):
        for j in range(i+1, len(txt_boxes)):
            a, b = txt_boxes[i], txt_boxes[j]
            ov = overlap(a, b)
            amin = min(a[2]*a[3], b[2]*b[3])
            if amin > 0 and ov / amin > 0.55 and ov > 0.25:
                print(f"[S{si}] TEXT OVERLAP {ov/amin*100:.0f}%: '{a[4][:22]}' ∩ '{b[4][:22]}'"); issues += 1

print(f"\n=== 레이아웃 이슈: {issues} ===\n")
print("=== 슬라이드별 텍스트 (내용 QA) ===")
for si, slide in enumerate(prs.slides, 1):
    txts = [b[4] for b in boxes(slide) if b[4]]
    print(f"\n--- S{si} ---")
    for t in txts:
        print("  " + t[:95].replace("\n", " / "))
